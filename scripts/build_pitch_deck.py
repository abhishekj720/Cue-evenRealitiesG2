"""Generate the Cue pitch deck as `docs/Cue-Pitch.pptx`.

Run:  .venv/bin/python scripts/build_pitch_deck.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "docs" / "Cue-Pitch.pptx"
SHOTS = Path(__file__).resolve().parent.parent / "docs" / "screenshots"

INK = RGBColor(0x0E, 0x0E, 0x0E)
BG = RGBColor(0xFA, 0xFA, 0xFA)
ACCENT = RGBColor(0x00, 0xB8, 0x6F)
MUTED = RGBColor(0x6B, 0x72, 0x80)
DARK = RGBColor(0x11, 0x18, 0x27)
HUD_BG = RGBColor(0x00, 0x00, 0x00)
HUD_FG = RGBColor(0x00, 0xE8, 0x3E)


def paint_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, *, size=28, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Helvetica"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return box


def add_bullets(slide, x, y, w, h, items, *, size=20, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "—  " + text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Helvetica"
        p.space_after = Pt(8)
    return box


def hud_card(slide, x, y, w, h, title, lines):
    """Mock G2 HUD card (black / green monospace)."""
    rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))  # MSO_SHAPE.RECTANGLE
    rect.fill.solid()
    rect.fill.fore_color.rgb = HUD_BG
    rect.line.color.rgb = HUD_FG
    rect.line.width = Pt(1.5)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(18); r.font.bold = True
    r.font.color.rgb = HUD_FG; r.font.name = "Menlo"
    for line in ["-" * 28, *lines]:
        pp = tf.add_paragraph()
        pp.alignment = PP_ALIGN.LEFT
        rr = pp.add_run(); rr.text = line
        rr.font.size = Pt(14); rr.font.color.rgb = HUD_FG
        rr.font.name = "Menlo"


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ---------- 1. Title ----------
    s = prs.slides.add_slide(blank_layout)
    paint_bg(s, DARK)
    add_text(s, 0.7, 2.5, 12, 1.5, "CUE", size=140, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, 0.7, 4.2, 12, 0.8, "Ambient social memory for the Even G2.", size=32, color=ACCENT)
    add_text(s, 0.7, 6.6, 12, 0.4, "Even Realities Builders' Day  ·  2026", size=14, color=MUTED)

    # ---------- 2. The hook ----------
    s = prs.slides.add_slide(blank_layout)
    paint_bg(s)
    add_text(s, 0.7, 0.6, 12, 0.5, "THE HOOK", size=12, bold=True, color=ACCENT)
    add_text(s, 0.7, 1.1, 12, 2.5,
             '"I met 40 people today.\nI remember maybe 6 names."',
             size=56, bold=True)
    add_text(s, 0.7, 5.3, 12, 0.8, "That's not a me problem — it's an interface problem.", size=24, color=MUTED)

    # ---------- 3. The wedge ----------
    s = prs.slides.add_slide(blank_layout)
    paint_bg(s)
    add_text(s, 0.7, 0.6, 12, 0.5, "THE WEDGE", size=12, bold=True, color=ACCENT)
    add_text(s, 0.7, 1.2, 12, 1.2,
             "Meta Ray-Bans remember faces.\nCue remembers people.",
             size=48, bold=True)
    # Two columns
    add_text(s, 0.7, 3.7, 5.5, 0.4, "Face-based (cameras)", size=16, bold=True, color=MUTED)
    add_bullets(s, 0.7, 4.1, 5.5, 3,
                ["Needs a camera on your face",
                 "Banned from bars, hospitals, therapy",
                 "Records strangers without consent"])
    add_text(s, 7, 3.7, 5.5, 0.4, "Voice-based (Cue)", size=16, bold=True, color=ACCENT)
    add_bullets(s, 7, 4.1, 5.5, 3,
                ["Cameraless by physics",
                 "Works everywhere cameras can't",
                 "Only remembers people you enrolled"])

    # ---------- 4. Live demo placeholder ----------
    s = prs.slides.add_slide(blank_layout)
    paint_bg(s, DARK)
    add_text(s, 0.7, 0.6, 12, 0.5, "LIVE DEMO", size=12, bold=True, color=ACCENT)
    add_text(s, 0.7, 2.7, 12, 1.5, "[ Walk up. Speak. Card appears. ]",
             size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_text(s, 0.7, 4.6, 12, 0.6,
             "Abhishek says one sentence.\nCue recognizes his voice. The HUD says: \"Ask about G2 BLE quirks.\"",
             size=20, color=ACCENT, align=PP_ALIGN.CENTER)

    # ---------- 5. What lands on the HUD (real screenshots) ----------
    s = prs.slides.add_slide(blank_layout)
    paint_bg(s)
    add_text(s, 0.7, 0.4, 12, 0.5, "WHAT LANDS ON THE HUD", size=12, bold=True, color=ACCENT)
    add_text(s, 0.7, 0.9, 12, 0.6, "One glanceable card. 576 x 288 px. Green on black.",
             size=18, color=MUTED)
    # Two real rendered cards side-by-side.
    for png, x in ((SHOTS / "03_match_abhishek@2x.png", 0.7),
                   (SHOTS / "05_match_priya@2x.png", 7.0)):
        if png.exists():
            s.shapes.add_picture(str(png), Inches(x), Inches(2.0),
                                 width=Inches(5.6))
    add_text(s, 0.7, 6.1, 12, 0.5,
             "Name  ·  what they do  ·  last thing promised  ·  what to say next",
             size=16, color=MUTED, align=PP_ALIGN.CENTER)

    # ---------- 5b. End-to-end HUD states montage ----------
    s = prs.slides.add_slide(blank_layout)
    paint_bg(s)
    add_text(s, 0.7, 0.4, 12, 0.5, "END-TO-END HUD STATES", size=12, bold=True, color=ACCENT)
    add_text(s, 0.7, 0.9, 12, 0.6, "Idle → match → another match → focus mode.",
             size=18, color=MUTED)
    positions = [
        (SHOTS / "01_idle@2x.png",           0.6, 1.8, "Idle"),
        (SHOTS / "03_match_abhishek@2x.png", 6.9, 1.8, "Match: Abhishek"),
        (SHOTS / "04_match_tim@2x.png",      0.6, 4.6, "Match: Tim"),
        (SHOTS / "06_focus@2x.png",          6.9, 4.6, "Focus mode"),
    ]
    for png, x, y, label in positions:
        if png.exists():
            s.shapes.add_picture(str(png), Inches(x), Inches(y), width=Inches(5.7))
            add_text(s, x, y + 2.55, 5.7, 0.3, label, size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # ---------- 5c. Push + pull ----------
    s = prs.slides.add_slide(blank_layout)
    paint_bg(s)
    add_text(s, 0.7, 0.4, 12, 0.5, "TWO WAYS TO RECALL", size=12, bold=True, color=ACCENT)
    add_text(s, 0.7, 0.9, 12, 0.8, "Not just a notification. A memory surface.", size=28, bold=True)
    # Two columns
    add_text(s, 0.7, 2.1, 5.8, 0.4, "Push  —  voice triggers it", size=18, bold=True, color=MUTED)
    add_bullets(s, 0.7, 2.6, 5.8, 3,
                ["Someone speaks near you",
                 "Cue matches their voiceprint",
                 "Card fades onto your HUD in ≤1.5s",
                 "You read, react, converse"], size=18)
    add_text(s, 7, 2.1, 5.8, 0.4, "Pull  —  you ask for it", size=18, bold=True, color=ACCENT)
    add_bullets(s, 7, 2.6, 5.8, 3,
                ["Triple-tap your right temple",
                 "Most-recent person fills the HUD",
                 "Cycle to older contacts with more taps",
                 "Before you walk into the next meeting"], size=18)
    add_text(s, 0.7, 6.2, 12, 0.5,
             "A memory you trust. Not a notification you react to.",
             size=18, color=MUTED, align=PP_ALIGN.CENTER)

    # ---------- 6. Architecture (flow diagram image) ----------
    s = prs.slides.add_slide(blank_layout)
    paint_bg(s)
    add_text(s, 0.7, 0.3, 12, 0.5, "HOW IT WORKS", size=12, bold=True, color=ACCENT)
    flow_png = SHOTS / "flow_diagram.png"
    if flow_png.exists():
        s.shapes.add_picture(str(flow_png), Inches(0.5), Inches(0.9),
                             width=Inches(12.3))

    # ---------- 6b. Terminal log ----------
    s = prs.slides.add_slide(blank_layout)
    paint_bg(s, DARK)
    add_text(s, 0.7, 0.4, 12, 0.5, "WHAT THE OPERATOR SEES", size=12, bold=True, color=ACCENT)
    add_text(s, 0.7, 0.9, 12, 0.6, "Terminal A — live recognition log on stage", size=20, color=RGBColor(0xEE, 0xEE, 0xEE))
    term_png = SHOTS / "07_terminal_log.png"
    if term_png.exists():
        s.shapes.add_picture(str(term_png), Inches(1.8), Inches(1.9),
                             width=Inches(9.7))

    # ---------- 7. Privacy ----------
    s = prs.slides.add_slide(blank_layout)
    paint_bg(s)
    add_text(s, 0.7, 0.6, 12, 0.5, "PRIVACY, BY DESIGN", size=12, bold=True, color=ACCENT)
    add_text(s, 0.7, 1.2, 12, 1.4,
             "Cue only remembers people\nyou chose to remember.",
             size=44, bold=True)
    add_bullets(s, 0.7, 4.3, 12, 3, [
        "Explicit enrollment — no row is written without a temple tap",
        "Local-first — ~/.cue/people.db lives on your phone. One-tap delete.",
        "Cameraless by physics — the G2 has no lens. No face database exists.",
        "No network on the critical path. Claude used only for optional summaries.",
    ], size=20)

    # ---------- 8. Why this resonates ----------
    s = prs.slides.add_slide(blank_layout)
    paint_bg(s)
    add_text(s, 0.7, 0.6, 12, 0.5, "WHY THIS RESONATES", size=12, bold=True, color=ACCENT)
    add_text(s, 0.7, 1.2, 12, 1.0, "Cue is the room.", size=44, bold=True)
    add_bullets(s, 0.7, 2.8, 12, 4, [
        "Everyone here forgot someone's name this week.",
        "You walked into 40 intros today. You'll forget 34.",
        "Meta tried cameras — they get you thrown out of the room.",
        "G2 stays in. Cue is Quiet Tech rendered as product.",
    ], size=22)

    # ---------- 9. What's next ----------
    s = prs.slides.add_slide(blank_layout)
    paint_bg(s)
    add_text(s, 0.7, 0.6, 12, 0.5, "WHAT'S NEXT", size=12, bold=True, color=ACCENT)
    add_text(s, 0.7, 1.1, 12, 1.0, "From hackathon to daily wear.", size=36, bold=True)
    # Two columns of roadmap cards.
    roadmap_left = [
        ("Speaker diarization",
         "Match multiple voices in group\nconversations, not just 1:1."),
        ("Proactive recall",
         '"You promised Priya a beta\ninvite 21 days ago." — surfaced\nwhen her voice is detected.'),
        ("Calendar hydration",
         "Pull last meeting notes for the\nmatched person from your cal."),
    ]
    roadmap_right = [
        ("Two-way consent",
         "Two G2 wearers enroll each\nother in one tap. Standardized\nconsent handshake."),
        ("CRM sync",
         "One-way export to HubSpot or\nNotion for sales workflows."),
        ("On-device Claude",
         "Brief generation runs on the\nphone via on-device models —\nzero cloud on the critical path."),
    ]

    def _card(slide, x, y, title, body):
        rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(5.9), Inches(1.4))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        rect.line.color.rgb = ACCENT
        rect.line.width = Pt(1.5)
        tf = rect.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.12)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.size = Pt(18); r.font.bold = True
        r.font.color.rgb = INK; r.font.name = "Helvetica"
        for line in body.split("\n"):
            pp = tf.add_paragraph()
            rr = pp.add_run(); rr.text = line
            rr.font.size = Pt(13); rr.font.color.rgb = MUTED
            rr.font.name = "Helvetica"

    for i, (t, b) in enumerate(roadmap_left):
        _card(s, 0.5, 2.2 + i * 1.6, t, b)
    for i, (t, b) in enumerate(roadmap_right):
        _card(s, 6.9, 2.2 + i * 1.6, t, b)

    # ---------- 10. Close ----------
    s = prs.slides.add_slide(blank_layout)
    paint_bg(s, DARK)
    add_text(s, 0.7, 1.5, 12, 1.2,
             "Faces are data.",
             size=56, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_text(s, 0.7, 2.8, 12, 1.2,
             "People are memory.",
             size=56, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, 0.7, 4.6, 12, 0.8,
             "G2 is the first wearable quiet enough\nto hold the second one.",
             size=22, color=RGBColor(0xDD, 0xDD, 0xDD), align=PP_ALIGN.CENTER)
    add_text(s, 0.7, 6.6, 12, 0.4,
             "Thank you.  ·  github.com/abhishekj720/Cue-evenRealitiesG2",
             size=14, color=MUTED, align=PP_ALIGN.CENTER)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT}  ({OUT.stat().st_size / 1024:.1f} KB)")


if __name__ == "__main__":
    build()
